import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation():
    try:
        prs = Presentation()
    except Exception as e:
        print(f"Error initializing presentation: {e}")
        return
        
    # Slide 1: Title Slide (Layout 0)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Project 15: Real-Time Security Log Visualizer"
    subtitle.text = "Cybersecurity Class Project\nImplementation and Demonstration"

    # Helper function to add content slides
    def add_slide(title_text, bullet_points):
        slide = prs.slides.add_slide(prs.slide_layouts[1]) # Layout 1 is Title and Content
        title = slide.shapes.title
        title.text = title_text
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
    
    # Slide 2: Problem and Background
    add_slide("Problem & Background", [
        "Increasing number of cyber threats require real-time visibility.",
        "System administrators need to detect malicious activities from raw logs promptly.",
        "Manual log parsing is tedious, error-prone, and too slow to stop live attacks.",
        "Need for a localized, lightweight, real-time threat detection tool."
    ])

    # Slide 3: Objective
    add_slide("Objective", [
        "Build a Real-Time Security Log Visualizer.",
        "Monitor a local 'access.log' file in real-time.",
        "Automatically parse newly appended log entries.",
        "Detect security threats: Brute-Force attacks and Off-Hour logins.",
        "Update a dynamic dashboard to provide immediate situational awareness."
    ])

    # Slide 4: Methodology and Approach
    add_slide("Methodology & Approach", [
        "Employed a 'File Watcher' methodology (Method 2).",
        "Asynchronous background process monitors target file.",
        "Triggers parsing routines instantly on file modification.",
        "Maintains historical state in-memory to evaluate patterns spanning multiple logs.",
        "Live mapping of aggregated data to visual dashboard interfaces."
    ])

    # Slide 5: Tools & Technologies Used
    add_slide("Tools & Technologies Used", [
        "Python 3: Core backend logic and data simulation.",
        "Streamlit: Rapid web framework for real-time interactive dashboards.",
        "Watchdog: Python library for monitoring file system events.",
        "Pandas: Data manipulation, filtering, and aggregation."
    ])

    # Slide 6: Implementation & Experiment Details
    add_slide("Implementation Details", [
        "Log Simulator (logger_sim.py): Generates mock success/failed logs, simulates brute-force behavior, and introduces off-hour login events.",
        "Dashboard Application (app.py): Reads logs and applies detection rules.",
        "Rule 1 - Brute Force: Flag IP address with > 5 'Failed Login' attempts.",
        "Rule 2 - Off-Hours: Flag logins outside the 9:00 AM - 6:00 PM window."
    ])

    # Slide 7: Results & Findings
    add_slide("Results & Findings", [
        "Log Simulator effectively streamed mock data.",
        "Dashboard updated in real-time with sub-second latency.",
        "Brute-force attacks were correctly identified and highlighted as 'Critical'.",
        "Off-hour logins were successfully tagged.",
        "Visual metrics provided clear, actionable insights."
    ])

    # Slide 8: Limitations
    add_slide("Limitations", [
        "Scalability: In-memory state tracking is not suitable for massive log volumes in production.",
        "Persistence: Application state resets on server restart without a database.",
        "Log Format rigidness: Hardcoded parsing logic breaks easily if the log format changes."
    ])

    # Slide 9: Lessons Learned
    add_slide("Lessons Learned", [
        "Integrating asynchronous event watchers (Watchdog) with synchronous UI frameworks (Streamlit) requires careful state management.",
        "Real-time visual feedback is vastly superior to static log analysis.",
        "Creating realistic simulated threat data is challenging but crucial for testing."
    ])

    # Slide 10: Recommendations for Future
    add_slide("Recommendations & Future Improvements", [
        "Database Integration: Connect a time-series database (e.g., InfluxDB) or Elasticsearch.",
        "Advanced Analytics: Implement Machine Learning algorithms to identify complex anomalies.",
        "Alerting Mechanism: Add Webhooks/Emails for automated admin notification on critical thresholds."
    ])

    output_path = "presentation.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation generated successfully at: {output_path}")
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == "__main__":
    create_presentation()
